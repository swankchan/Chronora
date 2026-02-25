#!/usr/bin/env python3
"""
樓盤廣告風格 PPTX：現代 template、安藤忠雄風格建築群，透光・現代感・地標性。
"""
import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("請先安裝: pip install python-pptx")
    raise

BASE_DIR = Path(__file__).resolve().parent
OUTPUTS = BASE_DIR / "outputs_base"
OUT_PPTX = BASE_DIR / "建築群設計介紹_安藤忠雄風格.pptx"

# 樓盤廣告風格配色（現代・高端）
COLOR_HEADER_BG = RGBColor(0x1a, 0x1a, 0x2e)   # 深藍黑
COLOR_ACCENT = RGBColor(0xc9, 0xa2, 0x27)      # 金
COLOR_WHITE = RGBColor(0xff, 0xff, 0xff)
COLOR_TITLE = RGBColor(0x1a, 0x1a, 0x2e)
COLOR_BODY = RGBColor(0x44, 0x44, 0x44)
COLOR_CAPTION = RGBColor(0x66, 0x66, 0x66)
HEADER_H = 0.42   # 頂欄高度（英寸）
FOOTER_H = 0.38   # 底欄高度（英寸）

def add_header_bar(slide, title_text=""):
    """每頁頂部：深色橫條 + 可選標題。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(HEADER_H)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_HEADER_BG
    shape.line.fill.background()
    if title_text:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.08), Inches(9), Inches(0.28))
        p = tb.text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

def add_footer_bar(slide, left_text="", right_text=""):
    """每頁底部：細橫條，可放案名 / 聯絡。"""
    footer_top = 7.5 - FOOTER_H
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(footer_top), Inches(10), Inches(FOOTER_H)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_HEADER_BG
    shape.line.fill.background()
    if left_text:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(footer_top + 0.06), Inches(5), Inches(0.25))
        tb.text_frame.paragraphs[0].text = left_text
        tb.text_frame.paragraphs[0].font.size = Pt(11)
        tb.text_frame.paragraphs[0].font.color.rgb = COLOR_ACCENT
    if right_text:
        tb = slide.shapes.add_textbox(Inches(5.5), Inches(footer_top + 0.06), Inches(4), Inches(0.25))
        tb.text_frame.paragraphs[0].text = right_text
        tb.text_frame.paragraphs[0].font.size = Pt(11)
        tb.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE

# 選幾張代表圖（依檔名排序取部分）
def get_sample_images(max_slides=6):
    if not OUTPUTS.exists():
        return []
    files = sorted(OUTPUTS.glob("chronora_base_*.png"))
    if len(files) <= max_slides:
        return [str(f) for f in files]
    step = max(1, len(files) // max_slides)
    return [str(files[i]) for i in range(0, len(files), step)][:max_slides]

# 拎齊所有圖（用於縮圖牆）
def get_all_images(max_count=100):
    if not OUTPUTS.exists():
        return []
    files = sorted(OUTPUTS.glob("chronora_base_*.png"))[:max_count]
    return [str(f) for f in files]

def add_slide_thumbnail_grid(prs, img_paths, title, cols=5, thumb_inches=1.42, margin=0.4, case_name=""):
    if not img_paths:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 先加縮圖（確保在內容層），再用頂/底欄疊在上面
    start_x = Inches(margin)
    start_y = Inches(HEADER_H + 0.35)
    gap = Inches(0.08)
    w = Inches(thumb_inches)
    h = Inches(thumb_inches)
    idx = 0
    row = 0
    while idx < len(img_paths):
        for col in range(cols):
            if idx >= len(img_paths):
                break
            x = start_x + col * (w + gap)
            y = start_y + row * (h + gap)
            bottom_inch = HEADER_H + 0.35 + (row + 1) * (thumb_inches + 0.08)
            if bottom_inch > 7.5 - FOOTER_H - 0.1:
                break
            path = os.path.abspath(img_paths[idx])
            if os.path.isfile(path):
                try:
                    slide.shapes.add_picture(path, x, y, width=w, height=h)
                except Exception as e:
                    print(f"縮圖無法加入 {path}: {e}")
            idx += 1
        row += 1
        if row >= 4:
            break
    add_header_bar(slide, title)
    add_footer_bar(slide, left_text=case_name, right_text="")

def add_cover_slide(prs, case_name, tagline, subtitle=""):
    """樓盤封面：頂欄 + 大案名 + 金標語 + 副標。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "  ")
    # 案名（大）
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(8.8), Inches(1.0))
    p = tx.text_frame.paragraphs[0]
    p.text = case_name
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    # 金標語
    tx2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.35), Inches(8.8), Inches(0.6))
    p2 = tx2.text_frame.paragraphs[0]
    p2.text = tagline
    p2.font.size = Pt(22)
    p2.font.color.rgb = COLOR_ACCENT
    if subtitle:
        tx3 = slide.shapes.add_textbox(Inches(0.6), Inches(4.0), Inches(8.8), Inches(0.5))
        tx3.text_frame.paragraphs[0].text = subtitle
        tx3.text_frame.paragraphs[0].font.size = Pt(16)
        tx3.text_frame.paragraphs[0].font.color.rgb = COLOR_BODY
    add_footer_bar(slide, left_text=case_name, right_text="敬請預約賞析")

def add_end_slide(prs, case_name):
    """樓盤結尾：感謝 + 預約賞析。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, case_name)
    tx = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(0.8))
    p = tx.text_frame.paragraphs[0]
    p.text = "感謝賞析"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    tx2 = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(0.6))
    tx2.text_frame.paragraphs[0].text = "歡迎預約參觀 · 親臨品鑑"
    tx2.text_frame.paragraphs[0].font.size = Pt(20)
    tx2.text_frame.paragraphs[0].font.color.rgb = COLOR_ACCENT
    add_footer_bar(slide, left_text=case_name, right_text="謝謝")

def add_slide_image_and_text(prs, title, bullets, img_path, bottom_thumbs=None, case_name=""):
    """左圖右文 + 頂/底欄；可選底部小縮圖。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, title)
    add_footer_bar(slide, left_text=case_name, right_text="")
    top = HEADER_H + 0.15
    main_h = Inches(4.0) if bottom_thumbs else Inches(4.6)
    slide.shapes.add_picture(img_path, Inches(0.4), Inches(top), width=Inches(4.2), height=main_h)
    tx = slide.shapes.add_textbox(Inches(4.8), Inches(top), Inches(4.7), Inches(0.6))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    y = top + 0.65
    for line in bullets:
        tb = slide.shapes.add_textbox(Inches(4.8), Inches(y), Inches(4.7), Inches(0.48))
        tb.text_frame.word_wrap = True
        tb.text_frame.paragraphs[0].text = line
        tb.text_frame.paragraphs[0].font.size = Pt(11)
        tb.text_frame.paragraphs[0].font.color.rgb = COLOR_BODY
        y += 0.46
    if bottom_thumbs:
        tw, th = Inches(0.54), Inches(0.54)
        gap = Inches(0.05)
        start_y = Inches(5.35)
        for i, path in enumerate(bottom_thumbs[:12]):
            try:
                x = Inches(0.4 + i * 0.59)
                if 0.4 + (i + 1) * 0.59 > 4.6:
                    break
                slide.shapes.add_picture(path, x, start_y, width=tw, height=th)
            except Exception:
                pass

def add_slide_one_image_and_caption(prs, img_path, title, caption_lines, case_name=""):
    """一頁一圖 + 頂/底欄 + 標題與說明。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, title)
    add_footer_bar(slide, left_text=case_name, right_text="")
    top = HEADER_H + 0.12
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(top), width=Inches(9), height=Inches(4.5))
    y = top + 4.6
    for line in caption_lines:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.35))
        tb.text_frame.word_wrap = True
        tb.text_frame.paragraphs[0].text = line
        tb.text_frame.paragraphs[0].font.size = Pt(12)
        tb.text_frame.paragraphs[0].font.color.rgb = COLOR_CAPTION
        y += 0.36

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    CASE_NAME = "光・影"
    TAGLINE = "安藤忠雄風格 · 現代地標建築"
    SUBTITLE = "清水混凝土與玻璃的對話 · 透光 · 現代感"

    sample_images = get_sample_images(8)
    all_images = get_all_images(100)
    thumbs_per_slide = 20

    # 1. 樓盤封面
    add_cover_slide(prs, CASE_NAME, TAGLINE, SUBTITLE)

    # 2. 建築理念（設計風格）
    if sample_images:
        thumbs_1 = [p for p in all_images if p != sample_images[0]][:12]
        add_slide_image_and_text(prs, "建築理念 · 向安藤忠雄致敬", [
            "• 清水混凝土與玻璃的對話：裸露混凝土搭配大面玻璃，簡潔而有力。",
            "• 極簡幾何與光影：乾淨線條、弧形與直線並置，強調光線與陰影。",
            "• 雕塑感量體：曲線屋頂、拱形入口等元素具辨識度。",
            "• 材質二元：玻璃的透與混凝土的實，虛實對比與現代感。",
            "• 秩序與留白：少即是多，結構與光線成為主角。",
        ], sample_images[0], bottom_thumbs=thumbs_1, case_name=CASE_NAME)

    # 3. 項目亮點（設計優點）
    if len(sample_images) >= 2:
        thumbs_2 = [p for p in all_images if p != sample_images[1]][:12]
        add_slide_image_and_text(prs, "項目亮點 · 生活質感", [
            "• 透光：大面玻璃引入自然光，室內明亮節能，隨時間呈現不同表情。",
            "• 現代感：流線造型、簡潔材質，傳達當代、專業的建築形象。",
            "• 地標性：獨特輪廓與曲面元素，在都市中具辨識度與記憶點。",
            "• 內外連結：透明立面拉近室內與城市、自然的關係，視野開闊。",
            "• 耐久與低維護：混凝土與玻璃的組合利於長期使用。",
        ], sample_images[1], bottom_thumbs=thumbs_2, case_name=CASE_NAME)

    # 4. 建築意象（大圖 + 一句話）
    captions = [
        ("透光與自然採光", ["大面玻璃引入自然光，室內明亮節能，隨時間與天氣呈現不同表情。"]),
        ("現代感與地標性", ["流線造型與簡潔材質，獨特輪廓在都市中具辨識度與記憶點。"]),
        ("光影與材質對話", ["清水混凝土與玻璃的虛實對比，光線與陰影成為空間主角。"]),
        ("內外連結與視野", ["透明立面拉近室內與城市、自然的關係，視野開闊。"]),
    ]
    for i, img_path in enumerate(sample_images[2:6] if len(sample_images) > 2 else []):
        title, lines = captions[i % len(captions)]
        add_slide_one_image_and_caption(prs, img_path, title, lines, case_name=CASE_NAME)

    # 5. 意象總覽（縮圖牆）
    for chunk_start in range(0, len(all_images), thumbs_per_slide):
        chunk = all_images[chunk_start : chunk_start + thumbs_per_slide]
        page = chunk_start // thumbs_per_slide + 1
        total_pages = (len(all_images) + thumbs_per_slide - 1) // thumbs_per_slide
        title = f"建築意象總覽（{page}/{total_pages}）"
        add_slide_thumbnail_grid(prs, chunk, title, cols=5, thumb_inches=1.42, case_name=CASE_NAME)

    # 6. 樓盤結尾
    add_end_slide(prs, CASE_NAME)

    prs.save(OUT_PPTX)
    print(f"已生成: {OUT_PPTX}")

if __name__ == "__main__":
    main()
